from pptx import Presentation
from typing import Dict


class PPTXParser:
    """A class for parsing text from PowerPoint presentations."""

    def __init__(self, path: str) -> None:
        """Initializes a new instance of the pptx_parser class.
        Args:
            path (str): The file path of the PowerPoint presentation to parse.
        Raises:
            FileNotFoundError: If the provided file path is not valid.
        """
        try:
            self.presentation = Presentation(path)
        except FileNotFoundError:
            raise FileNotFoundError("Please provide a valid pptx file path.")

        # dict[slide number : text extracted from slide]
        self.parsing_results_dict: Dict[int, str] = {}

    def extract_text_from_presentation(self) -> None:
        """
        Extracts the text from each slide in the presentation.
        Each slide's text is added to the parsing_result_dict attribute of the class.
        If no text was found in slide, empty str is placed in the dict.

       """
        for slide_number, slide in enumerate(self.presentation.slides):
            slide_text = PPTXParser.extract_text_from_slide(slide).replace('©', '')
            if slide_text:
                self.parsing_results_dict[slide_number + 1] = slide_text
            else:
                self.parsing_results_dict[slide_number + 1] = ''

    @staticmethod
    def extract_text_from_slide(current_slide) -> str:
        """
        Extracts the text from a single slide.
        Args:
            current_slide: The slide to extract text from.
        Returns:
            str: The extracted text from the slide.
        """
        return ''.join([PPTXParser.extract_text_from_shape(shape) for shape in current_slide.shapes])

    @staticmethod
    def extract_text_from_shape(shape):
        """
        Extracts the text, if exists, from a single shape in a slide.
        Args:
            shape: The shape to extract text from.
        Returns:
            str: The extracted text from the shape.
        """
        return ''.join(
            [PPTXParser.extract_text_from_paragraph(paragraph) for paragraph in
             shape.text_frame.paragraphs]) if hasattr(shape, 'text_frame') else ''

    @staticmethod
    def extract_text_from_paragraph(paragraph):
        """
        Extracts the text from a single paragraph in a shape.
        Args:
            paragraph: The paragraph to extract text from.
        Returns:
            str: The extracted text from the paragraph.
        """
        return ''.join([run.text + ' ' for run in paragraph.runs])


    def get_presentation_raw_text(self):
        """
        get all text from presentation concatenated into one single string.
        Returns:
        str: all the text of the presentation.
        """
        return ''.join([value + ' ' for value in self.parsing_results_dict.values()])

